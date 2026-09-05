import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Palette
    COLOR_PRIMARY = RGBColor(67, 56, 202)      # Indigo 700
    COLOR_DARK = RGBColor(15, 23, 42)          # Slate 900
    COLOR_TEXT = RGBColor(51, 65, 85)          # Slate 700
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_ACCENT = RGBColor(245, 158, 11)      # Amber 500
    COLOR_BG_CARD = RGBColor(248, 250, 252)    # Slate 50
    COLOR_BORDER = RGBColor(226, 232, 240)

    # -------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Dark Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK
    bg1.line.color.rgb = COLOR_DARK
    
    # Header tag
    tag = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(0.5))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "RAZORPAY BUILDATHON 2026 • TRACK 3"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    # Title
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(1.8))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI-Powered Revenue Recovery Agent"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Subtitle
    s_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.333), Inches(1.0))
    tf = s_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Autonomous, Intelligent & Compliant Recovery for Razorpay Merchants"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(148, 163, 184)
    
    # Presenter Card
    p_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.333), Inches(0.8))
    tf = p_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Presented by: Kashish Jindal"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 2: THE PROBLEM
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Slide Title
    title_box = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Problem: The Leaky Checkout Funnel"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 3 Cards for Problems
    cards_data = [
        ("10% – 15% Lost Revenue", "Merchants lose a massive share of sales to failed checkouts: UPI bank timeouts, card declines, and abandoned carts.", COLOR_PRIMARY),
        ("Generic & 'Dumb' Recovery", "Traditional recovery tools blast the exact same generic email to everyone, ignoring root cause and context.", COLOR_DARK),
        ("Spam & Compliance Risks", "Uncontrolled automated retries irritate customers, violate communication limits, and damage merchant brand reputation.", COLOR_ACCENT)
    ]
    
    left_positions = [Inches(1.0), Inches(5.0), Inches(9.0)]
    for i, (title, desc, color) in enumerate(cards_data):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[i], Inches(2.2), Inches(3.333), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1.5)
        
        tb = slide2.shapes.add_textbox(left_positions[i] + Inches(0.2), Inches(2.6), Inches(2.933), Inches(3.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = COLOR_TEXT

    # -------------------------------------------------------------
    # SLIDE 3: OUR SOLUTION
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    
    title_box = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Solution: Autonomous AI Revenue Recovery"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    sol_cards = [
        ("1. AI Root-Cause Diagnosis", "Evaluates exact Razorpay error codes (UPI timeout vs. limit decline) and determines the best channel & wait time.", COLOR_PRIMARY),
        ("2. Real Razorpay Integration", "Calls official Razorpay Orders API to dynamically generate secure checkout sessions matching the order amount.", COLOR_PRIMARY),
        ("3. Adaptive Hinglish Messaging", "Generates high-converting, personalized messages tailored for Indian phone numbers to win back buyers.", COLOR_PRIMARY),
        ("4. Strict Compliance Guardrails", "Terminates after 3 attempts and stops instantly if a customer replies 'STOP' to guarantee zero spam.", COLOR_PRIMARY)
    ]
    
    grid_pos = [
        (Inches(1.0), Inches(2.2)),
        (Inches(7.0), Inches(2.2)),
        (Inches(1.0), Inches(4.5)),
        (Inches(7.0), Inches(4.5))
    ]
    
    for i, (title, desc, color) in enumerate(sol_cards):
        x, y = grid_pos[i]
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.333), Inches(1.9))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1.5)
        
        tb = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(4.933), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_TEXT

    # -------------------------------------------------------------
    # SLIDE 4: ARCHITECTURE & DEMO
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    
    title_box = slide4.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How It Works: The 4-Step Recovery Loop"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    steps = [
        ("Step 1", "Payment Fails", "Razorpay sends webhook event (or triggered via simulator)."),
        ("Step 2", "AI Diagnoses", "Agent determines root cause, channel & personalized copy."),
        ("Step 3", "Checkout Created", "Real Razorpay Order & Checkout link generated dynamically."),
        ("Step 4", "Payment Saved", "Customer pays -> callback marks revenue recovered in real-time.")
    ]
    
    step_x = [Inches(1.0), Inches(4.0), Inches(7.0), Inches(10.0)]
    for i, (num, heading, detail) in enumerate(steps):
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, step_x[i], Inches(2.2), Inches(2.5), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_CARD
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1.5)
        
        tb = slide4.shapes.add_textbox(step_x[i] + Inches(0.15), Inches(2.5), Inches(2.2), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        p2 = tf.add_paragraph()
        p2.text = heading
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_PRIMARY
        
        p3 = tf.add_paragraph()
        p3.text = "\n" + detail
        p3.font.size = Pt(14)
        p3.font.color.rgb = COLOR_TEXT

    # -------------------------------------------------------------
    # SLIDE 5: CONCLUSION & IMPACT
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = COLOR_DARK
    bg5.line.color.rgb = COLOR_DARK
    
    t_box = slide5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(1.2))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Recovering Margin. Respecting Customers."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    summary_box = slide5.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(2.0))
    tf = summary_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "✓ Measurable batch recoveries tracked in real-time"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(203, 213, 225)
    
    p2 = tf.add_paragraph()
    p2.text = "✓ Powered by Google Gemini & Official Razorpay SDK"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    
    p3 = tf.add_paragraph()
    p3.text = "✓ Full audit trail and compliance stopping rules"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    
    git_box = slide5.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.333), Inches(0.8))
    tf = git_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GitHub: https://github.com/KashishJindal20/razorpay-revenue-recovery"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    output_path = "presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
